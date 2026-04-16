import os
import shutil
import zipfile
import tempfile
import subprocess
import uuid
from typing import Any, List, Optional, Dict
from fastapi import APIRouter, UploadFile, File, HTTPException
from pydantic import BaseModel, Field
import aiohttp
import asyncio
import xml.etree.ElementTree as ET
import re

from pptx import Presentation

from services.documents_loader import DocumentsLoader
from utils.asset_directory_utils import get_images_directory
import uuid
from constants.documents import POWERPOINT_TYPES


PPTX_SLIDES_ROUTER = APIRouter(prefix="/pptx-slides", tags=["PPTX Slides"])
DESIGN_CANVAS_WIDTH = 1280
DESIGN_CANVAS_HEIGHT = 720


class SlideData(BaseModel):
    slide_number: int
    screenshot_url: str
    xml_content: str
    normalized_fonts: List[str]


class FontAnalysisResult(BaseModel):
    internally_supported_fonts: List[
        Dict[str, str]
    ]  # [{"name": "Open Sans", "google_fonts_url": "..."}]
    not_supported_fonts: List[str]  # ["Custom Font Name"]


class PptxSlidesResponse(BaseModel):
    success: bool
    slides: List[SlideData]
    total_slides: int
    fonts: Optional[FontAnalysisResult] = None


class PptxTextBlock(BaseModel):
    text: str
    kind: str
    left: Optional[int] = None
    top: Optional[int] = None
    width: Optional[int] = None
    height: Optional[int] = None
    canvas_left: Optional[float] = None
    canvas_top: Optional[float] = None
    canvas_width: Optional[float] = None
    canvas_height: Optional[float] = None
    source_order: int = 0
    table_rows: List[List[str]] = Field(default_factory=list)
    bullets: List[str] = Field(default_factory=list)
    numeric_series: List[Dict[str, Any]] = Field(default_factory=list)


class PptxSlideText(BaseModel):
    slide_number: int
    text: str
    xml_text: str
    notes: Optional[str] = None
    blocks: List[PptxTextBlock]
    slide_width: int
    slide_height: int
    slide_aspect: str
    canvas_width: int = DESIGN_CANVAS_WIDTH
    canvas_height: int = DESIGN_CANVAS_HEIGHT
    bullets: List[str] = Field(default_factory=list)
    numeric_series: List[Dict[str, Any]] = Field(default_factory=list)


class PptxTextExtractionResponse(BaseModel):
    success: bool
    slides: List[PptxSlideText]
    total_slides: int
    deck_width: int
    deck_height: int
    deck_aspect: str
    canvas_width: int = DESIGN_CANVAS_WIDTH
    canvas_height: int = DESIGN_CANVAS_HEIGHT


# NEW: Fonts-only router and response for PPTX
class PptxFontsResponse(BaseModel):
    success: bool
    fonts: FontAnalysisResult


PPTX_FONTS_ROUTER = APIRouter(prefix="/pptx-fonts", tags=["PPTX Fonts"])

# NEW: Normalize font family names by removing style/weight/stretch descriptors and splitting camel case
_STYLE_TOKENS = {
    # styles
    "italic",
    "italics",
    "ital",
    "oblique",
    "roman",
    # combined style shortcuts
    "bolditalic",
    "bolditalics",
    # weights
    "thin",
    "hairline",
    "extralight",
    "ultralight",
    "light",
    "demilight",
    "semilight",
    "book",
    "regular",
    "normal",
    "medium",
    "semibold",
    "demibold",
    "bold",
    "extrabold",
    "ultrabold",
    "black",
    "extrablack",
    "ultrablack",
    "heavy",
    # width/stretch
    "narrow",
    "condensed",
    "semicondensed",
    "extracondensed",
    "ultracondensed",
    "expanded",
    "semiexpanded",
    "extraexpanded",
    "ultraexpanded",
}
# Modifiers commonly used with style tokens
_STYLE_MODIFIERS = {"semi", "demi", "extra", "ultra"}


def _insert_spaces_in_camel_case(value: str) -> str:
    # Insert space before capital letters preceded by lowercase or digits (e.g., MontserratBold -> Montserrat Bold)
    value = re.sub(r"(?<=[a-z0-9])([A-Z])", r" \1", value)
    # Handle sequences like BoldItalic -> Bold Italic
    value = re.sub(r"([A-Z]+)([A-Z][a-z])", r"\1 \2", value)
    return value


def normalize_font_family_name(raw_name: str) -> str:
    if not raw_name:
        return raw_name
    # Replace separators with spaces
    name = raw_name.replace("_", " ").replace("-", " ")
    # Insert spaces in camel case
    name = _insert_spaces_in_camel_case(name)
    # Collapse multiple spaces
    name = re.sub(r"\s+", " ", name).strip()
    # Lowercase helper for matching but keep original casing for output
    lower_name = name.lower()
    # Quick cut: if the full string ends with a pure style suffix, trim it
    for style in sorted(_STYLE_TOKENS, key=len, reverse=True):
        if lower_name.endswith(" " + style):
            name = name[: -(len(style) + 1)]
            lower_name = lower_name[: -(len(style) + 1)]
            break
    # Tokenize
    tokens_original = name.split(" ")
    tokens_filtered: List[str] = []
    for index, tok in enumerate(tokens_original):
        lower_tok = tok.lower()
        # Always keep the first token to avoid stripping families like "Black Ops One"
        if index == 0:
            tokens_filtered.append(tok)
            continue
        # Drop style tokens and standalone modifiers
        if lower_tok in _STYLE_TOKENS or lower_tok in _STYLE_MODIFIERS:
            continue
        tokens_filtered.append(tok)
    # If everything except first token was dropped and first token is a style token (unlikely), fallback to original
    if not tokens_filtered:
        tokens_filtered = tokens_original
    normalized = " ".join(tokens_filtered).strip()
    # Final cleanup of leftover multiple spaces
    normalized = re.sub(r"\s+", " ", normalized)
    return normalized


def extract_fonts_from_oxml(xml_content: str) -> List[str]:
    """
    Extract font names from OXML content.

    Args:
        xml_content: OXML content as string

    Returns:
        List of unique font names found in the OXML
    """
    fonts = set()

    try:
        # Parse the XML content
        root = ET.fromstring(xml_content)

        # Define namespaces commonly used in OXML
        namespaces = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }

        # Search for font references in various OXML elements
        # Look for latin fonts
        for font_elem in root.findall(".//a:latin", namespaces):
            if "typeface" in font_elem.attrib:
                fonts.add(font_elem.attrib["typeface"])

        # Look for east asian fonts
        for font_elem in root.findall(".//a:ea", namespaces):
            if "typeface" in font_elem.attrib:
                fonts.add(font_elem.attrib["typeface"])

        # Look for complex script fonts
        for font_elem in root.findall(".//a:cs", namespaces):
            if "typeface" in font_elem.attrib:
                fonts.add(font_elem.attrib["typeface"])

        # Look for font references in theme elements
        for font_elem in root.findall(".//a:font", namespaces):
            if "typeface" in font_elem.attrib:
                fonts.add(font_elem.attrib["typeface"])

        # Look for rPr (run properties) font references
        for rpr_elem in root.findall(".//a:rPr", namespaces):
            for font_elem in rpr_elem.findall(".//a:latin", namespaces):
                if "typeface" in font_elem.attrib:
                    fonts.add(font_elem.attrib["typeface"])

        # Also search without namespace prefix for compatibility
        for font_elem in root.findall(".//latin"):
            if "typeface" in font_elem.attrib:
                fonts.add(font_elem.attrib["typeface"])

        # Regex fallback for fonts that might be missed
        font_pattern = r'typeface="([^"]+)"'
        regex_fonts = re.findall(font_pattern, xml_content)
        fonts.update(regex_fonts)

        # Filter out system fonts and empty values
        system_fonts = {"+mn-lt", "+mj-lt", "+mn-ea", "+mj-ea", "+mn-cs", "+mj-cs", ""}
        fonts = {font for font in fonts if font not in system_fonts and font.strip()}

        return list(fonts)

    except Exception as e:
        print(f"Error extracting fonts from OXML: {e}")
        return []


async def check_google_font_availability(font_name: str) -> bool:
    """
    Check if a font is available in Google Fonts.

    Args:
        font_name: Name of the font to check

    Returns:
        True if font is available in Google Fonts, False otherwise
    """
    try:
        formatted_name = font_name.replace(" ", "+")
        url = f"https://fonts.googleapis.com/css2?family={formatted_name}&display=swap"

        async with aiohttp.ClientSession() as session:
            async with session.head(
                url, timeout=aiohttp.ClientTimeout(total=10)
            ) as response:
                return response.status == 200

    except Exception as e:
        print(f"Error checking Google Font availability for {font_name}: {e}")
        return False


async def analyze_fonts_in_all_slides(slide_xmls: List[str]) -> FontAnalysisResult:
    """
    Analyze fonts across all slides and determine Google Fonts availability.

    Args:
        slide_xmls: List of OXML content strings from all slides

    Returns:
        FontAnalysisResult with supported and unsupported fonts
    """
    # Extract fonts from all slides
    raw_fonts = set()
    for xml_content in slide_xmls:
        slide_fonts = extract_fonts_from_oxml(xml_content)
        raw_fonts.update(slide_fonts)

    # Normalize to root families (e.g., "Montserrat Italic" -> "Montserrat")
    normalized_fonts = {normalize_font_family_name(f) for f in raw_fonts}
    # Remove empties if any
    normalized_fonts = {f for f in normalized_fonts if f}

    if not normalized_fonts:
        return FontAnalysisResult(internally_supported_fonts=[], not_supported_fonts=[])

    # Check each normalized font's availability in Google Fonts concurrently
    tasks = [check_google_font_availability(font) for font in normalized_fonts]
    results = await asyncio.gather(*tasks)

    internally_supported_fonts = []
    not_supported_fonts = []

    for font, is_available in zip(normalized_fonts, results):
        if is_available:
            formatted_name = font.replace(" ", "+")
            google_fonts_url = f"https://fonts.googleapis.com/css2?family={formatted_name}&display=swap"
            internally_supported_fonts.append(
                {"name": font, "google_fonts_url": google_fonts_url}
            )
        else:
            not_supported_fonts.append(font)

    return FontAnalysisResult(
        internally_supported_fonts=internally_supported_fonts, not_supported_fonts=[]
    )


def _is_valid_pptx_upload(upload_file: UploadFile) -> bool:
    if upload_file.content_type in POWERPOINT_TYPES:
        return True
    return bool(upload_file.filename and upload_file.filename.lower().endswith(".pptx"))


MOJIBAKE_MARKERS = (
    "Ð",
    "Ñ",
    "Â",
    "â",
    "Р°",
    "Рµ",
    "Рё",
    "Рѕ",
    "Рґ",
    "Р»",
    "Рј",
    "РЅ",
    "Рї",
    "СЃ",
    "С‚",
    "СЊ",
    "СЌ",
    "вЂ",
)


def _mojibake_score(value: str) -> int:
    return sum(value.count(marker) for marker in MOJIBAKE_MARKERS)


def _cyrillic_score(value: str) -> int:
    return len(re.findall(r"[\u0400-\u04FF]", value or ""))


def _fix_mojibake_text(value: str) -> str:
    if not value or _mojibake_score(value) == 0:
        return value

    candidates = [value]
    for encoding in ("latin1", "cp1251"):
        try:
            candidates.append(value.encode(encoding).decode("utf-8"))
        except Exception:
            continue

    def score(candidate: str) -> tuple[int, int]:
        return (_cyrillic_score(candidate) - _mojibake_score(candidate) * 4, -_mojibake_score(candidate))

    return max(candidates, key=score)


def _clean_extracted_text(value: str) -> str:
    value = _fix_mojibake_text(value or "")
    value = value.replace("\u00a0", " ")
    value = re.sub(r"[ \t]+\n", "\n", value)
    value = re.sub(r"\n{4,}", "\n\n\n", value)
    return value.strip()


def _parse_numeric_value(raw_value: str) -> Optional[float]:
    normalized = raw_value.replace(" ", "").replace("\u00a0", "")
    if "," in normalized and "." in normalized:
        normalized = normalized.replace(",", "")
    else:
        normalized = normalized.replace(",", ".")
    normalized = re.sub(r"[^0-9.+-]", "", normalized)
    if not normalized or normalized in {"+", "-", ".", "+.", "-."}:
        return None
    try:
        return float(normalized)
    except ValueError:
        return None


def _extract_bullets(text: str) -> List[str]:
    bullets = []
    for line in _clean_extracted_text(text).splitlines():
        match = re.match(r"^\s*(?:\d{1,2}[\).]\s+|[-*\u2022]\s+)(.+)$", line)
        if match:
            bullets.append(match.group(1).strip())
    return bullets


def _extract_numeric_series(text: str) -> List[Dict[str, Any]]:
    series: List[Dict[str, Any]] = []
    for line in _clean_extracted_text(text).splitlines():
        if "|" in line or "\t" in line:
            continue
        match = re.search(r"[-+]?\d[\d\s.,]*(?:%|[A-Za-zА-Яа-я]{1,8})?", line)
        if not match:
            continue
        raw = match.group(0).strip()
        value = _parse_numeric_value(raw)
        if value is None:
            continue
        suffix_match = re.search(r"(%|[A-Za-zА-Яа-я]{1,8})$", raw)
        suffix = suffix_match.group(1) if suffix_match else ""
        label = (line[: match.start()] + line[match.end() :]).strip(" :-\t")
        series.append(
            {
                "label": label or line.strip(),
                "raw": raw,
                "value": value,
                "suffix": suffix,
            }
        )
    return series[:12]


def _slide_aspect(width: int, height: int) -> str:
    if not width or not height:
        return "16:9"
    ratio = width / height
    if abs(ratio - (16 / 9)) < 0.03:
        return "16:9"
    if abs(ratio - (4 / 3)) < 0.03:
        return "4:3"
    return f"{ratio:.3f}:1"


def _shape_position(shape, slide_width: int, slide_height: int) -> dict:
    left = int(shape.left) if getattr(shape, "left", None) is not None else None
    top = int(shape.top) if getattr(shape, "top", None) is not None else None
    width = int(shape.width) if getattr(shape, "width", None) is not None else None
    height = int(shape.height) if getattr(shape, "height", None) is not None else None

    def scale(value: Optional[int], source: int, target: int) -> Optional[float]:
        if value is None or not source:
            return None
        return round((value / source) * target, 2)

    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "canvas_left": scale(left, slide_width, DESIGN_CANVAS_WIDTH),
        "canvas_top": scale(top, slide_height, DESIGN_CANVAS_HEIGHT),
        "canvas_width": scale(width, slide_width, DESIGN_CANVAS_WIDTH),
        "canvas_height": scale(height, slide_height, DESIGN_CANVAS_HEIGHT),
    }


def _paragraph_text(paragraph) -> str:
    run_text = "".join(run.text for run in paragraph.runs)
    return _clean_extracted_text(run_text if run_text else paragraph.text)


def _text_frame_text(text_frame) -> str:
    paragraphs = [_paragraph_text(paragraph) for paragraph in text_frame.paragraphs]
    return _clean_extracted_text("\n".join(paragraph for paragraph in paragraphs if paragraph))


def _extract_text_blocks_from_shapes(
    shapes,
    slide_width: int,
    slide_height: int,
) -> List[PptxTextBlock]:
    blocks: List[PptxTextBlock] = []

    for shape in sorted(
        shapes,
        key=lambda item: (
            int(getattr(item, "top", 0) or 0),
            int(getattr(item, "left", 0) or 0),
        ),
    ):
        if hasattr(shape, "shapes"):
            blocks.extend(
                _extract_text_blocks_from_shapes(
                    shape.shapes,
                    slide_width,
                    slide_height,
                )
            )
            continue

        if getattr(shape, "has_table", False):
            rows = []
            for row in shape.table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = _text_frame_text(cell.text_frame)
                    if cell_text:
                        cells.append(cell_text)
                if cells:
                    rows.append(cells)
            table_text = _clean_extracted_text("\n".join(" | ".join(row) for row in rows))
            if table_text:
                blocks.append(
                    PptxTextBlock(
                        text=table_text,
                        kind="table",
                        table_rows=rows,
                        bullets=_extract_bullets(table_text),
                        numeric_series=_extract_numeric_series(table_text),
                        **_shape_position(shape, slide_width, slide_height),
                    )
                )
            continue

        if getattr(shape, "has_text_frame", False):
            text = _text_frame_text(shape.text_frame)
            if text:
                blocks.append(
                    PptxTextBlock(
                        text=text,
                        kind="text",
                        bullets=_extract_bullets(text),
                        numeric_series=_extract_numeric_series(text),
                        **_shape_position(shape, slide_width, slide_height),
                    )
                )

    return blocks


def _extract_xml_text(xml_content: str) -> str:
    try:
        root = ET.fromstring(xml_content)
        namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        values = [node.text for node in root.findall(".//a:t", namespaces) if node.text]
        return _clean_extracted_text("\n".join(values))
    except Exception:
        return ""


def _extract_notes_text(slide) -> Optional[str]:
    try:
        if not slide.has_notes_slide:
            return None
        notes_frame = slide.notes_slide.notes_text_frame
        if not notes_frame:
            return None
        return _text_frame_text(notes_frame) or None
    except Exception:
        return None


def _merge_missing_xml_text_blocks(
    blocks: List[PptxTextBlock], xml_text: str
) -> List[PptxTextBlock]:
    if not xml_text:
        return blocks

    block_text = "\n".join(block.text for block in blocks)
    missing_lines = [
        line
        for line in xml_text.splitlines()
        if line.strip() and line.strip() not in block_text
    ]

    if missing_lines:
        missing_text = _clean_extracted_text("\n".join(missing_lines))
        blocks.append(
            PptxTextBlock(
                text=missing_text,
                kind="xml_fallback",
                bullets=_extract_bullets(missing_text),
                numeric_series=_extract_numeric_series(missing_text),
            )
        )

    return blocks


def _with_source_order(blocks: List[PptxTextBlock]) -> List[PptxTextBlock]:
    ordered_blocks = sorted(
        blocks,
        key=lambda item: (
            item.canvas_top if item.canvas_top is not None else 999999,
            item.canvas_left if item.canvas_left is not None else 999999,
            item.source_order,
        ),
    )
    for index, block in enumerate(ordered_blocks):
        block.source_order = index
    return ordered_blocks


@PPTX_SLIDES_ROUTER.post("/extract-text", response_model=PptxTextExtractionResponse)
async def extract_pptx_text(
    pptx_file: UploadFile = File(..., description="PPTX file to extract text from"),
    include_notes: bool = False,
):
    """
    Extract all readable text from a PPTX file, grouped by slide.

    The main extraction uses python-pptx so text boxes and tables keep slide-level
    grouping and approximate visual order. The raw slide XML is also scanned for
    <a:t> nodes and appended as an xml_fallback block when text is present in XML
    but not found through python-pptx.
    """

    if not _is_valid_pptx_upload(pptx_file):
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type. Expected PPTX file, got {pptx_file.content_type}",
        )

    if (
        hasattr(pptx_file, "size")
        and pptx_file.size
        and pptx_file.size > (100 * 1024 * 1024)
    ):
        raise HTTPException(
            status_code=400,
            detail="PPTX file exceeded max upload size of 100 MB",
        )

    with tempfile.TemporaryDirectory() as temp_dir:
        pptx_path = os.path.join(temp_dir, "presentation.pptx")
        with open(pptx_path, "wb") as f:
            f.write(await pptx_file.read())

        try:
            presentation = Presentation(pptx_path)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to read PPTX file: {e}",
            )

        slide_xmls = _extract_slide_xmls(pptx_path, temp_dir)
        slides_text: List[PptxSlideText] = []
        deck_width = int(presentation.slide_width)
        deck_height = int(presentation.slide_height)
        deck_aspect = _slide_aspect(deck_width, deck_height)

        for index, slide in enumerate(presentation.slides):
            xml_text = _extract_xml_text(slide_xmls[index]) if index < len(slide_xmls) else ""
            blocks = _extract_text_blocks_from_shapes(
                slide.shapes,
                deck_width,
                deck_height,
            )
            blocks = _with_source_order(_merge_missing_xml_text_blocks(blocks, xml_text))
            slide_text = _clean_extracted_text(
                "\n\n".join(block.text for block in blocks if block.text)
            )
            notes = _extract_notes_text(slide) if include_notes else None
            bullets = []
            numeric_series = []
            for block in blocks:
                bullets.extend(block.bullets)
                numeric_series.extend(block.numeric_series)

            slides_text.append(
                PptxSlideText(
                    slide_number=index + 1,
                    text=slide_text,
                    xml_text=xml_text,
                    notes=notes,
                    blocks=blocks,
                    slide_width=deck_width,
                    slide_height=deck_height,
                    slide_aspect=deck_aspect,
                    bullets=bullets[:24],
                    numeric_series=numeric_series[:24],
                )
            )

        return PptxTextExtractionResponse(
            success=True,
            slides=slides_text,
            total_slides=len(slides_text),
            deck_width=deck_width,
            deck_height=deck_height,
            deck_aspect=deck_aspect,
        )


@PPTX_SLIDES_ROUTER.post("/process", response_model=PptxSlidesResponse)
async def process_pptx_slides(
    pptx_file: UploadFile = File(..., description="PPTX file to process"),
    fonts: Optional[List[UploadFile]] = File(None, description="Optional font files"),
):
    """
    Process a PPTX file to extract slide screenshots and XML content.

    This endpoint:
    1. Validates the uploaded PPTX file
    2. Installs any provided font files
    3. Unzips the PPTX to extract slide XMLs
    4. Uses LibreOffice to generate slide screenshots
    5. Returns both screenshot URLs and XML content for each slide
    """

    # Validate PPTX file
    if pptx_file.content_type not in POWERPOINT_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type. Expected PPTX file, got {pptx_file.content_type}",
        )
    # Enforce 100MB size limit
    if (
        hasattr(pptx_file, "size")
        and pptx_file.size
        and pptx_file.size > (100 * 1024 * 1024)
    ):
        raise HTTPException(
            status_code=400,
            detail="PPTX file exceeded max upload size of 100 MB",
        )

    # Create temporary directory for processing
    with tempfile.TemporaryDirectory() as temp_dir:
        if True:
            # Save uploaded PPTX file
            pptx_path = os.path.join(temp_dir, "presentation.pptx")
            with open(pptx_path, "wb") as f:
                pptx_content = await pptx_file.read()
                f.write(pptx_content)

            # Install fonts if provided
            if fonts:
                await _install_fonts(fonts, temp_dir)

            # Extract slide XMLs from PPTX
            slide_xmls = _extract_slide_xmls(pptx_path, temp_dir)

            # Convert PPTX to PDF
            pdf_path = await _convert_pptx_to_pdf(pptx_path, temp_dir)

            # Generate screenshots using LibreOffice
            screenshot_paths = await DocumentsLoader.get_page_images_from_pdf_async(
                pdf_path, temp_dir
            )
            print(f"Screenshot paths: {screenshot_paths}")

            # Analyze fonts across all slides
            font_analysis = await analyze_fonts_in_all_slides(slide_xmls)
            print(
                f"Font analysis completed: {len(font_analysis.internally_supported_fonts)} supported, {len(font_analysis.not_supported_fonts)} not supported"
            )

            # Move screenshots to images directory and generate URLs
            images_dir = get_images_directory()
            presentation_id = uuid.uuid4()
            presentation_images_dir = os.path.join(images_dir, str(presentation_id))
            os.makedirs(presentation_images_dir, exist_ok=True)

            slides_data = []

            for i, (xml_content, screenshot_path) in enumerate(
                zip(slide_xmls, screenshot_paths), 1
            ):
                # Move screenshot to permanent location
                screenshot_filename = f"slide_{i}.png"
                permanent_screenshot_path = os.path.join(
                    presentation_images_dir, screenshot_filename
                )

                if (
                    os.path.exists(screenshot_path)
                    and os.path.getsize(screenshot_path) > 0
                ):
                    # Use shutil.copy2 instead of os.rename to handle cross-device moves
                    shutil.copy2(screenshot_path, permanent_screenshot_path)
                    screenshot_url = (
                        f"/app_data/images/{presentation_id}/{screenshot_filename}"
                    )
                else:
                    # Fallback if screenshot generation failed or file is empty placeholder
                    screenshot_url = "/static/images/placeholder.jpg"

                # Compute normalized fonts for this slide
                raw_slide_fonts = extract_fonts_from_oxml(xml_content)
                normalized_fonts = sorted(
                    {normalize_font_family_name(f) for f in raw_slide_fonts if f}
                )

                slides_data.append(
                    SlideData(
                        slide_number=i,
                        screenshot_url=screenshot_url,
                        xml_content=xml_content,
                        normalized_fonts=normalized_fonts,
                    )
                )

            return PptxSlidesResponse(
                success=True,
                slides=slides_data,
                total_slides=len(slides_data),
                fonts=font_analysis,
            )


# NEW: Fonts-only endpoint leveraging the same font extraction/analysis
@PPTX_FONTS_ROUTER.post("/process", response_model=PptxFontsResponse)
async def process_pptx_fonts(
    pptx_file: UploadFile = File(..., description="PPTX file to analyze fonts from")
):
    """
    Analyze a PPTX file and return only the fonts used in the document.

    Uses the exact same font extraction and analysis utilities as the /pptx-slides endpoint.
    """
    # Validate PPTX file
    if pptx_file.content_type not in POWERPOINT_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type. Expected PPTX file, got {pptx_file.content_type}",
        )

    # Create temporary directory for processing
    with tempfile.TemporaryDirectory() as temp_dir:
        # Save uploaded PPTX file
        pptx_path = os.path.join(temp_dir, "presentation.pptx")
        with open(pptx_path, "wb") as f:
            pptx_content = await pptx_file.read()
            f.write(pptx_content)

        # Extract slide XMLs from PPTX
        slide_xmls = _extract_slide_xmls(pptx_path, temp_dir)

        # Analyze fonts across all slides (same logic as in /pptx-slides)
        font_analysis = await analyze_fonts_in_all_slides(slide_xmls)

        return PptxFontsResponse(
            success=True,
            fonts=font_analysis,
        )


def _create_font_alias_config(raw_fonts: List[str]) -> str:
    """Create a temporary fontconfig configuration that aliases variant family names to normalized root families.
    Returns the path to the config file.
    """
    # Build mapping from raw -> normalized where different
    mappings: Dict[str, str] = {}
    for f in raw_fonts:
        normalized = normalize_font_family_name(f)
        if normalized and normalized != f:
            mappings[f] = normalized
    # Create config only if we have mappings
    fd, fonts_conf_path = tempfile.mkstemp(prefix="fonts_alias_", suffix=".conf")
    os.close(fd)
    with open(fonts_conf_path, "w", encoding="utf-8") as cfg:
        cfg.write(
            """<?xml version='1.0'?>
<!DOCTYPE fontconfig SYSTEM "urn:fontconfig:fonts.dtd">
<fontconfig>
  <include>/etc/fonts/fonts.conf</include>
"""
        )
        for src, dst in mappings.items():
            cfg.write(
                f"""
  <match target="pattern">
    <test name="family" compare="eq">
      <string>{src}</string>
    </test>
    <edit name="family" mode="assign" binding="strong">
      <string>{dst}</string>
    </edit>
  </match>
"""
            )
        cfg.write("\n</fontconfig>\n")
    return fonts_conf_path


async def _install_fonts(fonts: List[UploadFile], temp_dir: str) -> None:
    """Install provided font files to the system."""
    fonts_dir = os.path.join(temp_dir, "fonts")
    os.makedirs(fonts_dir, exist_ok=True)

    for font_file in fonts:
        # Save font file
        font_path = os.path.join(fonts_dir, font_file.filename)
        with open(font_path, "wb") as f:
            font_content = await font_file.read()
            f.write(font_content)

        # Install font (copy to system fonts directory)
        try:
            subprocess.run(
                ["cp", font_path, "/usr/share/fonts/truetype/"],
                check=True,
                capture_output=True,
            )
        except subprocess.CalledProcessError as e:
            print(f"Warning: Failed to install font {font_file.filename}: {e}")

    # Refresh font cache
    try:
        subprocess.run(["fc-cache", "-f", "-v"], check=True, capture_output=True)
    except subprocess.CalledProcessError as e:
        print(f"Warning: Failed to refresh font cache: {e}")


def _extract_slide_xmls(pptx_path: str, temp_dir: str) -> List[str]:
    """Extract slide XML content from PPTX file."""
    slide_xmls = []
    extract_dir = os.path.join(temp_dir, "pptx_extract")

    try:
        # Unzip PPTX file
        with zipfile.ZipFile(pptx_path, "r") as zip_ref:
            zip_ref.extractall(extract_dir)

        # Look for slides in ppt/slides/ directory
        slides_dir = os.path.join(extract_dir, "ppt", "slides")

        if not os.path.exists(slides_dir):
            raise Exception("No slides directory found in PPTX file")

        # Get all slide XML files and sort them numerically
        slide_files = [
            f
            for f in os.listdir(slides_dir)
            if f.startswith("slide") and f.endswith(".xml")
        ]
        slide_files.sort(key=lambda x: int(x.replace("slide", "").replace(".xml", "")))

        # Read XML content from each slide
        for slide_file in slide_files:
            slide_path = os.path.join(slides_dir, slide_file)
            with open(slide_path, "r", encoding="utf-8") as f:
                slide_xmls.append(f.read())

        return slide_xmls

    except Exception as e:
        raise Exception(f"Failed to extract slide XMLs: {str(e)}")


async def _convert_pptx_to_pdf(pptx_path: str, temp_dir: str) -> str:
    """Generate PNG screenshots of PPTX slides using LibreOffice + ImageMagick."""
    screenshots_dir = os.path.join(temp_dir, "screenshots")
    os.makedirs(screenshots_dir, exist_ok=True)

    try:
        # First, get the number of slides by extracting XMLs
        slide_xmls = _extract_slide_xmls(pptx_path, temp_dir)
        slide_count = len(slide_xmls)

        # Build font alias config to force variant families to resolve to normalized root families
        raw_fonts: List[str] = []
        for xml in slide_xmls:
            raw_fonts.extend(extract_fonts_from_oxml(xml))
        raw_fonts = list({f for f in raw_fonts if f})
        fonts_conf_path = _create_font_alias_config(raw_fonts)
        env = os.environ.copy()
        env["FONTCONFIG_FILE"] = fonts_conf_path

        print(f"Found {slide_count} slides in presentation")

        # Step 1: Convert PPTX to PDF using LibreOffice
        print("Starting LibreOffice PDF conversion...")
        pdf_filename = "temp_presentation.pdf"
        pdf_path = os.path.join(screenshots_dir, pdf_filename)

        try:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    screenshots_dir,
                    pptx_path,
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=500,
                env=env,
            )

            print(f"LibreOffice PDF conversion output: {result.stdout}")
            if result.stderr:
                print(f"LibreOffice PDF conversion warnings: {result.stderr}")
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice PDF conversion timed out after 120 seconds")
        except subprocess.CalledProcessError as e:
            error_msg = e.stderr if e.stderr else str(e)
            raise Exception(f"LibreOffice PDF conversion failed: {error_msg}")

        # Find the generated PDF file (LibreOffice uses original filename)
        pdf_files = [f for f in os.listdir(screenshots_dir) if f.endswith(".pdf")]
        if not pdf_files:
            raise Exception("LibreOffice failed to generate PDF file")

        actual_pdf_path = os.path.join(screenshots_dir, pdf_files[0])
        print(f"Generated PDF: {actual_pdf_path}")
        return actual_pdf_path

    except Exception as e:
        # Re-raise the specific exceptions we've already handled
        if "timed out" in str(e) or "failed:" in str(e):
            raise
        # Handle any other unexpected exceptions
        raise Exception(f"Screenshot generation failed: {str(e)}")
